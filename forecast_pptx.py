"""
Generador de la presentación PPTX "Línea Base Forecast" (Ingeniería y Equipo Móvil).

La portada y la slide de cierre se clonan 1:1 desde el formato corporativo de JPV
(assets/plantilla_jpv.pptx, provisto por el usuario) — mismo fondo, logos y estilo,
solo cambia el texto dinámico (título y fecha). El título de las demás slides usa la
misma tipografía del formato (Arial Black, naranja de marca) que la slide interior de
ese documento. El resto del contenido usa una paleta clara con foco en una idea por
slide, buen contraste y tamaños de fuente legibles en pantalla grande.

Fase 2: genera 15 slides (Portada, Agenda, Título sección, Línea Base Real, Proyección
y Pipeline, Casos en Proceso Administrativo [si aplica], Forecast Total, Análisis de
Desviación, Evolución 7+5, Facturación Mensual, 3 slides de Top 10 casos por Pérdida
Bruta segmentados por Probabilidad de Cierre [100%, 75% y menor a 50%], y Cierre). Las
2 slides comparativas que dependen del historial acumulado de forecasts (Comparativo
Forecasts y Tabla Comparativa) llegan en la Fase 3, junto con la persistencia del
historial.
"""
import copy
import io as _io_clon
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# --- Paleta clara con acentos de marca JPV (naranja de la plantilla corporativa) ---
COLOR_BG = "FFFFFF"
COLOR_BAR = "1A365D"
COLOR_CARD_BG = "F5F7FA"
COLOR_CARD_BORDER = "D9E0E8"
COLOR_IE = "2B6CB0"
COLOR_EM = "C05621"
COLOR_CTM = "00796B"
COLOR_VERDE = "2F855A"
COLOR_ROJO = "C53030"
COLOR_TEXTO = "1A202C"
COLOR_MUTED = "718096"
COLOR_MUTED2 = "4A5568"
COLOR_TITULO = "FF9933"
FUENTE = "Calibri"
FUENTE_TITULO = "Arial Black"

# Tamaño de slide 13.333in x 7.5in (16:9 estándar), igual al formato corporativo de JPV.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGEN_X = Emu(609600)
CONTENIDO_W = Emu(10972800)

_ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
LOGO_PATH = os.path.join(_ASSETS_DIR, "logo_jpv_navy.png")
LOGO_ASPECT = 1351 / 208  # ancho/alto del archivo real
PLANTILLA_PATH = os.path.join(_ASSETS_DIR, "plantilla_jpv.pptx")


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _sin_borde(shape):
    shape.line.fill.background()


def nueva_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = _rgb(COLOR_BG)
    _sin_borde(fondo)
    fondo.shadow.inherit = False
    return slide


def barra(slide, y):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Emu(97536))
    b.fill.solid()
    b.fill.fore_color.rgb = _rgb(COLOR_BAR)
    _sin_borde(b)
    b.shadow.inherit = False
    return b


def caja(slide, x, y, w, h, fill_hex=COLOR_CARD_BG, border_hex=COLOR_CARD_BORDER):
    c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = _rgb(fill_hex)
    if border_hex:
        c.line.color.rgb = _rgb(border_hex)
        c.line.width = Pt(1)
    else:
        _sin_borde(c)
    c.shadow.inherit = False
    return c


def texto(slide, x, y, w, h, contenido, size, color_hex, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margenes=True, wrap=True, fuente=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    if not margenes:
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lineas = contenido if isinstance(contenido, list) else [contenido]
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = linea
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = fuente or FUENTE
        run.font.color.rgb = _rgb(color_hex)
    return box


def pie_de_pagina(slide, fuente_texto):
    """Footer estándar de las slides de contenido: texto de fuente a la izquierda y logo a la derecha."""
    texto(slide, MARGEN_X, Emu(6278880), Emu(10363200), Emu(219456),
          fuente_texto, 9, COLOR_MUTED2, margenes=False)
    alto_logo = Emu(219456)
    ancho_logo = Emu(int(alto_logo * LOGO_ASPECT))
    x_logo = Emu(11338560 - ancho_logo)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, x_logo, Emu(6278880), height=alto_logo)


def encabezado_slide(slide, titulo, subtitulo, color_subtitulo=COLOR_IE):
    texto(slide, MARGEN_X, Emu(219456), CONTENIDO_W, Emu(512064), titulo, 24, COLOR_TITULO,
          bold=True, fuente=FUENTE_TITULO)
    texto(slide, MARGEN_X, Emu(731520), CONTENIDO_W, Emu(304800), subtitulo, 14, color_subtitulo, bold=True)


def _fmt_uf(valor):
    return f"{valor:,.0f}".replace(",", ".")


def _fmt_uf2(valor):
    return f"{valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


# ---------------------------------------------------------
# Clonado de shapes desde la plantilla corporativa JPV (portada y cierre)
# ---------------------------------------------------------
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_HYPERLINK = _R_NS + "/hyperlink"


def _remapear_relaciones_clon(el, part_origen, part_destino):
    """Reasigna los r:embed/r:id de un elemento clonado a nuevas relaciones válidas en el
    part de destino: las imágenes se copian tal cual, los hipervínculos externos se recrean."""
    attr_embed, attr_link, attr_id = (f"{{{_R_NS}}}embed", f"{{{_R_NS}}}link", f"{{{_R_NS}}}id")
    for nodo in el.iter():
        for attr in (attr_embed, attr_link):
            rid = nodo.get(attr)
            if not rid:
                continue
            rel_origen = part_origen.rels.get(rid)
            if rel_origen is None or rel_origen.is_external:
                continue
            _, nuevo_rid = part_destino.get_or_add_image_part(
                _io_clon.BytesIO(rel_origen.target_part.blob)
            )
            nodo.set(attr, nuevo_rid)
        rid = nodo.get(attr_id)
        if not rid or rid not in part_origen.rels:
            continue
        rel_origen = part_origen.rels[rid]
        if rel_origen.reltype == _REL_HYPERLINK and rel_origen.is_external:
            nuevo_rid = part_destino.relate_to(rel_origen.target_ref, rel_origen.reltype, is_external=True)
            nodo.set(attr_id, nuevo_rid)


def clonar_shapes_de_plantilla(slide_destino, slide_origen):
    """Copia todos los shapes (imágenes, grupos, formas libres, textos) de slide_origen a
    slide_destino tal cual, remapeando relaciones de imagen/hipervínculo. Se usa para
    reproducir 1:1 la portada y el cierre del formato corporativo de JPV."""
    tree_destino = slide_destino.shapes._spTree
    part_origen = slide_origen.part
    part_destino = slide_destino.part
    tags_shape = {qn("p:sp"), qn("p:pic"), qn("p:grpSp"), qn("p:graphicFrame"), qn("p:cxnSp")}
    for el in list(slide_origen.shapes._spTree):
        if el.tag not in tags_shape:
            continue
        nuevo_el = copy.deepcopy(el)
        _remapear_relaciones_clon(nuevo_el, part_origen, part_destino)
        tree_destino.append(nuevo_el)


def _cargar_plantilla_jpv():
    return Presentation(PLANTILLA_PATH) if os.path.exists(PLANTILLA_PATH) else None


def _reemplazar_texto_clonado(slide, nombre_shape, lineas, size=None):
    """Reemplaza el contenido de un shape de texto clonado (por nombre), conservando el
    formato (fuente/color/negrita) del primer run original."""
    shape = next((s for s in slide.shapes if s.name == nombre_shape), None)
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    r0 = p0.runs[0] if p0.runs else None
    fuente_run = r0.font.name if r0 and r0.font.name else FUENTE_TITULO
    negrita = r0.font.bold if r0 else True
    size_run = size or (r0.font.size if r0 and r0.font.size else Pt(24))
    try:
        color = r0.font.color.rgb if r0 and r0.font.color and r0.font.color.type is not None else None
    except Exception:
        color = None
    if color is None:
        color = _rgb("FFFFFF")

    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for run in list(p0.runs):
        run._r.getparent().remove(run._r)

    lineas = lineas if isinstance(lineas, list) else [lineas]
    for i, linea in enumerate(lineas):
        p = p0 if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = linea
        run.font.name = fuente_run
        run.font.size = size_run
        run.font.bold = negrita
        run.font.color.rgb = color


# ---------------------------------------------------------
# SLIDE 1 — PORTADA (clonada del formato corporativo JPV)
# ---------------------------------------------------------
def slide_portada(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    plantilla = _cargar_plantilla_jpv()
    if plantilla is not None:
        clonar_shapes_de_plantilla(slide, plantilla.slides[0])
        _reemplazar_texto_clonado(slide, "CuadroTexto 9", [
            "Línea Base Forecast",
            f"Ingeniería y Equipo Móvil {datos['anio']}",
            f"Forecast {datos['label']} — {datos['nombre_mes_corte']} {datos['anio']}",
        ])
        _reemplazar_texto_clonado(slide, "CuadroTexto 1", datos["fecha_emision"])
    else:
        # Respaldo si falta el asset de la plantilla (assets/plantilla_jpv.pptx).
        fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        fondo.fill.solid()
        fondo.fill.fore_color.rgb = _rgb(COLOR_BG)
        _sin_borde(fondo)
        fondo.shadow.inherit = False
        barra(slide, 0)
        barra(slide, Emu(6339840))
        texto(slide, Emu(853440), Emu(1706880), Emu(10485120), Emu(853440),
              "Ingeniería y Equipo Móvil", 42, COLOR_EM, bold=True, fuente=FUENTE_TITULO)
        texto(slide, Emu(853440), Emu(2560320), Emu(10485120), Emu(548640),
              f"Línea Base {datos['anio']}: Forecast {datos['label']} — {datos['nombre_mes_corte']} {datos['anio']}",
              20, COLOR_TEXTO)
        texto(slide, Emu(853440), Emu(6217920), Emu(6096000), Emu(243840),
              datos["fecha_emision"], 10, COLOR_MUTED2)
    return slide


# SLIDE 2 — AGENDA
# ---------------------------------------------------------
def slide_agenda(prs, datos):
    slide = nueva_slide(prs)
    texto(slide, MARGEN_X, Emu(609600), CONTENIDO_W, Emu(609600), "Agenda", 36, COLOR_TITULO,
          bold=True, fuente=FUENTE_TITULO)

    items = [
        f"Forecast {datos['label']} {datos['anio']}",
        "Análisis de Desempeño Histórico",
        "La Mirada de la Industria",
        "Plan de Acción",
    ]
    y = 1706880
    for i, item in enumerate(items, start=1):
        caja(slide, MARGEN_X, Emu(y), Emu(731520), Emu(731520), fill_hex=COLOR_CARD_BG, border_hex=None)
        texto(slide, MARGEN_X, Emu(y), Emu(731520), Emu(731520), str(i), 29.5, COLOR_CTM, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(457200 + 640080), Emu(y), Emu(8534400), Emu(731520), item, 21.5, COLOR_TEXTO,
              margenes=False)
        y += 853440

    pie_de_pagina(slide, "JPV Asociados | sedgwick.")
    return slide


# ---------------------------------------------------------
# SLIDE 3 — TÍTULO DE SECCIÓN
# ---------------------------------------------------------
def slide_titulo_seccion(prs, datos):
    slide = nueva_slide(prs)
    texto(slide, Emu(853440), Emu(2560320), Emu(10485120), Emu(1280160),
          [f"Forecast Financiero {datos['anio']}: Análisis de", "Producción y Facturación."],
          36, COLOR_TITULO, bold=True, fuente=FUENTE_TITULO)
    texto(slide, Emu(853440), Emu(3840480), Emu(10485120), Emu(548640),
          "La Brecha del Forecast (La Realidad Incómoda)", 18.5, COLOR_MUTED)
    pie_de_pagina(slide, "JPV | sedgwick.")
    return slide


# ---------------------------------------------------------
# SLIDE 4 — LÍNEA BASE REAL
# ---------------------------------------------------------
def slide_linea_base_real(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(
        slide,
        "1. Línea Base Real (Cierre Facturación)",
        f"Estado Actual (YTD {datos['nombre_mes_inicio']} a {datos['nombre_mes_corte']})",
        COLOR_IE,
    )

    tarjetas = [
        ("Equipo Móvil", _fmt_uf2(datos['em_ytd']) + " UF", f"{datos['nombre_mes_inicio']}–{datos['nombre_mes_corte']} {datos['anio']}"),
        ("Ingeniería y Energía", _fmt_uf2(datos['ie_ytd']) + " UF", f"{datos['nombre_mes_inicio']}–{datos['nombre_mes_corte']} {datos['anio']}"),
        ("Total Consolidado YTD", _fmt_uf2(datos['ytd_total']) + " UF", f"{datos['nombre_mes_inicio']}–{datos['nombre_mes_corte']} {datos['anio']}"),
    ]
    x = 426720
    for titulo_t, valor_t, sub_t in tarjetas:
        caja(slide, Emu(x), Emu(1219200), Emu(2987040), Emu(1280160))
        texto(slide, Emu(x), Emu(1292352), Emu(2987040), Emu(268224), titulo_t, 12.0, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(1560576), Emu(2987040), Emu(487680), valor_t, 26.5, COLOR_TEXTO, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(2048256), Emu(2987040), Emu(280416), sub_t, 10.5, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        x += 2987040 + 109728

    x = 426720
    tarjetas2 = [
        ("Honorarios en Stock", _fmt_uf(datos['pipeline_bruto_total']) + " UF", "Pipeline bruto total " + str(datos['anio'])),
        ("Meta Anual", _fmt_uf(datos['meta']) + " UF", str(datos['anio'])),
    ]
    for titulo_t, valor_t, sub_t in tarjetas2:
        caja(slide, Emu(x), Emu(2682240), Emu(2987040), Emu(1280160))
        texto(slide, Emu(x), Emu(2755392), Emu(2987040), Emu(268224), titulo_t, 12.0, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(3023616), Emu(2987040), Emu(487680), valor_t, 26.5, COLOR_TEXTO, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(3511296), Emu(2987040), Emu(280416), sub_t, 10.5, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        x += 2987040 + 109728

    texto(slide, MARGEN_X, Emu(4267200), Emu(5333333), Emu(304800), "Análisis de contexto", 14.0, COLOR_IE, bold=True)
    bullets = [f"• {b}" for b in datos.get('bullets_linea_base', [])]
    if bullets:
        texto(slide, MARGEN_X, Emu(4632960), Emu(10972800), Emu(609600), bullets, 12.5, COLOR_TEXTO, margenes=False)

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE 5 — PROYECCIÓN Y PIPELINE
# ---------------------------------------------------------
def slide_proyeccion_pipeline(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(
        slide,
        f"2. Proyección y Pipeline {datos['label']} (Resto del año)",
        f"Proyección (Forecast {datos['nombre_mes_corte']}—Diciembre)",
        COLOR_IE,
    )

    proy_em = datos['proy_em']
    caja(slide, Emu(426720), Emu(1219200), Emu(5333333), Emu(1706880))
    texto(slide, Emu(548640), Emu(1292352), Emu(5066667), Emu(304800),
          "Equipo Móvil (Pipeline Hon. Probables)", 13.5, COLOR_EM, bold=True, margenes=False)
    texto(slide, Emu(548640), Emu(1645920), Emu(5066667), Emu(426720),
          f"{_fmt_uf(proy_em['em_stock'])} UF (stock)  ·  {_fmt_uf(proy_em['em_promedio_total'])} UF (promedio × meses)",
          17.5, COLOR_TEXTO, bold=True, margenes=False)
    meses_txt = "-".join(m[:3] for m in proy_em['meses_usados']) or "—"
    texto(slide, Emu(548640), Emu(2133600), Emu(5066667), Emu(609600),
          f"Complemento EM — Promedio últimos {len(proy_em['meses_usados'])} meses\n"
          f"{_fmt_uf(proy_em['prom_em_3m'])} UF/mes ({meses_txt}) × {datos['meses_proyectados']} meses = {_fmt_uf(proy_em['em_promedio_total'])} UF",
          12.0, COLOR_MUTED, margenes=False)

    if datos['total_admin'] > 0:
        caja(slide, Emu(5943600), Emu(1219200), Emu(5402400), Emu(1706880), border_hex=COLOR_CTM)
        texto(slide, Emu(6065520), Emu(1292352), Emu(5066667), Emu(304800),
              "Proceso Administrativo de Facturación", 13.5, COLOR_CTM, bold=True, margenes=False)
        texto(slide, Emu(6065520), Emu(1645920), Emu(5066667), Emu(426720),
              f"{_fmt_uf(datos['total_admin'])} UF", 24.0, COLOR_TEXTO, bold=True, margenes=False)
        texto(slide, Emu(6065520), Emu(2133600), Emu(5066667), Emu(609600),
              f"{len(datos['casos_admin'])} caso(s) en proceso administrativo",
              12.0, COLOR_MUTED, margenes=False)

    proy_ie = datos['proy_ie']
    texto(slide, MARGEN_X, Emu(3169920), Emu(8000000), Emu(304800),
          f"Ingeniería — Proyección {datos['nombre_mes_corte']}—Diciembre:", 14.0, COLOR_IE, bold=True)

    ie_tarjetas = [
        ("Casos menores", "(<1.000 UF pérdida)", _fmt_uf(proy_ie['ie_menores_proj']) + " UF",
         f"{_fmt_uf(proy_ie['prom_ie_menores'])} UF/mes × {datos['meses_proyectados']} meses"),
        ("Casos mayores pipeline", "(≥1.000 UF, prob.)", _fmt_uf(proy_ie['ie_mayores_proj']) + " UF",
         f"{proy_ie['n_casos_mayores']} casos · Hon. Probables ponderados"),
        ("Engie/CTM", "(pipeline, prob.)", _fmt_uf(proy_ie['ie_engie_proj']) + " UF",
         f"{proy_ie['n_casos_engie']} caso(s) · Hon. Probables ponderados"),
    ]
    x = 426720
    for t1, t2, valor_t, sub_t in ie_tarjetas:
        caja(slide, Emu(x), Emu(3657600), Emu(3474720), Emu(1463040))
        texto(slide, Emu(x), Emu(3730752), Emu(3474720), Emu(240000), t1, 12.5, COLOR_MUTED, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(3970752), Emu(3474720), Emu(213333), t2, 10.5, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(4200000), Emu(3474720), Emu(426720), valor_t, 22.5, COLOR_TEXTO, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(4693333), Emu(3474720), Emu(373333), sub_t, 10.5, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        x += 3474720 + 73333

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE — CASOS EN PROCESO ADMINISTRATIVO (condicional)
# ---------------------------------------------------------
def slide_casos_admin(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(slide, "Casos en Proceso Administrativo de Facturación",
                      f"Total: {_fmt_uf(datos['total_admin'])} UF", COLOR_CTM)

    y = 1463040
    for caso in datos['casos_admin']:
        caja(slide, MARGEN_X, Emu(y), CONTENIDO_W, Emu(609600))
        texto(slide, Emu(457200 + 137160), Emu(y), Emu(8000000), Emu(609600),
              caso['nombre'] or "(sin nombre)", 16.0, COLOR_TEXTO, margenes=False)
        texto(slide, Emu(9066667), Emu(y), Emu(1829333), Emu(609600),
              f"{_fmt_uf(caso['monto'])} UF", 17.5, COLOR_CTM, bold=True,
              align=PP_ALIGN.RIGHT, margenes=False)
        y += 670560
        if y > 5600000:
            break

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE — FORECAST TOTAL
# ---------------------------------------------------------
def slide_forecast_total(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(slide, f"2. Forecast {datos['label']} (Total año {datos['anio']})",
                      f"Proyección (Forecast {datos['nombre_mes_corte']}—Diciembre)", COLOR_IE)

    tarjetas = [
        ("Equipo Móvil", COLOR_EM, _fmt_uf(datos['em_total']) + " UF",
         f"{_fmt_uf2(datos['em_ytd'])} YTD + {_fmt_uf(datos['em_proyectado'])} proyectado"),
        ("Ingeniería", COLOR_IE, _fmt_uf(datos['ie_total']) + " UF",
         f"{_fmt_uf2(datos['ie_ytd'])} YTD + {_fmt_uf(datos['ie_proyectado'])} proyectado"),
    ]
    x = 426720
    for titulo_t, color_t, valor_t, sub_t in tarjetas:
        caja(slide, Emu(x), Emu(1463040), Emu(3719407), Emu(1706880))
        texto(slide, Emu(x + 91440), Emu(1536192), Emu(3466667), Emu(304800), titulo_t, 14.0, color_t, bold=True, margenes=False)
        texto(slide, Emu(x + 91440), Emu(1889760), Emu(3466667), Emu(487680), valor_t, 29.5, COLOR_TEXTO, bold=True, margenes=False)
        texto(slide, Emu(x + 91440), Emu(2377440), Emu(3466667), Emu(426720), sub_t, 12.0, COLOR_MUTED, margenes=False)
        x += 3719407 + 121927

    if datos['total_admin'] > 0:
        caja(slide, Emu(x), Emu(1463040), Emu(3719407), Emu(1706880), border_hex=COLOR_CTM)
        texto(slide, Emu(x + 91440), Emu(1536192), Emu(3466667), Emu(304800),
              "Proceso Administrativo", 14.0, COLOR_CTM, bold=True, margenes=False)
        texto(slide, Emu(x + 91440), Emu(1889760), Emu(3466667), Emu(487680),
              _fmt_uf(datos['total_admin']) + " UF", 29.5, COLOR_TEXTO, bold=True, margenes=False)
        texto(slide, Emu(x + 91440), Emu(2377440), Emu(3466667), Emu(426720),
              "Cierre esperado próximas semanas", 12.0, COLOR_MUTED, margenes=False)

    texto(slide, MARGEN_X, Emu(3474720), Emu(8000000), Emu(304800),
          "Cierre de año proyectado global", 16.0, COLOR_TEXTO, bold=True)

    color_gap_sin = COLOR_VERDE if datos['cum_sin'] >= 100 else COLOR_EM
    caja(slide, Emu(426720), Emu(3962400), Emu(5242560), Emu(1584960))
    texto(slide, Emu(426720), Emu(4084320), Emu(5242560), Emu(487680),
          _fmt_uf(datos['cierre_sin_admin']) + " UF", 29.5, COLOR_TEXTO, bold=True,
          align=PP_ALIGN.CENTER, margenes=False)
    texto(slide, Emu(426720), Emu(4572000), Emu(5242560), Emu(304800),
          f"Sin proceso administrativo · {datos['cum_sin']:.1f}% de la meta", 12.5, color_gap_sin,
          align=PP_ALIGN.CENTER, margenes=False)

    if datos['total_admin'] > 0:
        color_gap_con = COLOR_VERDE if datos['cum_con'] >= 100 else COLOR_EM
        caja(slide, Emu(5791200), Emu(3962400), Emu(5242560), Emu(1584960), border_hex=COLOR_CTM)
        texto(slide, Emu(5791200), Emu(4084320), Emu(5242560), Emu(487680),
              _fmt_uf(datos['cierre_con_admin']) + " UF", 29.5, COLOR_TEXTO, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(5791200), Emu(4572000), Emu(5242560), Emu(304800),
              f"Con proceso administrativo · {datos['cum_con']:.1f}% de la meta", 12.5, color_gap_con,
              align=PP_ALIGN.CENTER, margenes=False)

    texto(slide, MARGEN_X, Emu(5730240), CONTENIDO_W, Emu(304800),
          f"Honorarios en Stock: {_fmt_uf(datos['pipeline_bruto_total'])} UF no facturados en pipeline total (bruto, sin probabilidad)",
          12.0, COLOR_MUTED)

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE — ANÁLISIS DE DESVIACIÓN
# ---------------------------------------------------------
def slide_desviacion(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(slide, f"3. Análisis de Desviación (Gap) vs. Meta — Forecast {datos['label']}",
                      "Evaluación de Brecha Estratégica y Acciones Mitigadoras", COLOR_IE)

    usar_con_admin = datos['total_admin'] > 0
    cierre_ref = datos['cierre_con_admin'] if usar_con_admin else datos['cierre_sin_admin']
    gap_ref = datos['gap_con'] if usar_con_admin else datos['gap_sin']
    cum_ref = datos['cum_con'] if usar_con_admin else datos['cum_sin']
    color_gap = COLOR_VERDE if gap_ref >= 0 else COLOR_ROJO

    etiquetas = [
        ("Meta Oficial Total", _fmt_uf(datos['meta']) + " UF", str(datos['anio']), COLOR_TEXTO),
        ("Cierre Proyectado", _fmt_uf(cierre_ref) + " UF",
         "Con proc. admin." if usar_con_admin else "Sin proc. admin.", COLOR_TEXTO),
        ("Desviación (Gap)", f"{gap_ref:+,.0f}".replace(",", ".") + " UF",
         f"Cumplimiento: {cum_ref:.1f}%", color_gap),
    ]
    x = 426720
    for titulo_t, valor_t, sub_t, color_t in etiquetas:
        caja(slide, Emu(x), Emu(1341120), Emu(3474720), Emu(1219200))
        texto(slide, Emu(x), Emu(1414272), Emu(3474720), Emu(268224), titulo_t, 12.0, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(1682496), Emu(3474720), Emu(426720), valor_t, 24.0, color_t, bold=True,
              align=PP_ALIGN.CENTER, margenes=False)
        texto(slide, Emu(x), Emu(2133600), Emu(3474720), Emu(280416), sub_t, 11.5, COLOR_MUTED,
              align=PP_ALIGN.CENTER, margenes=False)
        x += 3474720 + 73333

    texto(slide, Emu(426720), Emu(2804160), Emu(5242560), Emu(304800),
          "Acción Requerida: Equipo Móvil", 14.0, COLOR_EM, bold=True)
    bullets_em = [f"• {b}" for b in datos.get('bullets_accion_em', [])]
    texto(slide, Emu(426720), Emu(3169920), Emu(5242560), Emu(2438400), bullets_em, 12.0, COLOR_TEXTO, margenes=False)

    texto(slide, Emu(5943600), Emu(2804160), Emu(5242560), Emu(304800),
          "Acción Requerida: Ingeniería", 14.0, COLOR_IE, bold=True)
    bullets_ie = [f"• {b}" for b in datos.get('bullets_accion_ie', [])]
    texto(slide, Emu(5943600), Emu(3169920), Emu(5242560), Emu(2438400), bullets_ie, 12.0, COLOR_TEXTO, margenes=False)

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# Helpers de estilo para gráficos nativos (tema oscuro)
# ---------------------------------------------------------
def _estilizar_chart(chart, mostrar_leyenda=True):
    chart.has_title = False
    if mostrar_leyenda:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = _rgb(COLOR_MUTED)
        chart.legend.font.name = FUENTE
    else:
        chart.has_legend = False

    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(12)
    cat_ax.tick_labels.font.color.rgb = _rgb(COLOR_MUTED)
    cat_ax.tick_labels.font.name = FUENTE
    cat_ax.format.line.color.rgb = _rgb(COLOR_CARD_BORDER)

    val_ax = chart.value_axis
    val_ax.tick_labels.font.size = Pt(12)
    val_ax.tick_labels.font.color.rgb = _rgb(COLOR_MUTED)
    val_ax.tick_labels.font.name = FUENTE
    val_ax.format.line.color.rgb = _rgb(COLOR_CARD_BORDER)
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = _rgb(COLOR_CARD_BORDER)
    val_ax.major_gridlines.format.line.width = Pt(0.5)


def _colorear_series(chart, colores):
    for i, serie in enumerate(chart.series):
        color = colores[i % len(colores)]
        if chart.chart_type in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
            serie.format.line.color.rgb = _rgb(color)
            serie.format.line.width = Pt(2.5)
            serie.smooth = False
        else:
            serie.format.fill.solid()
            serie.format.fill.fore_color.rgb = _rgb(color)


# ---------------------------------------------------------
# SLIDE — EVOLUCIÓN N+M (gráfico de curvas acumuladas)
# ---------------------------------------------------------
def slide_evolucion(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(slide, f"Evolución Forecast {datos['label']} año {datos['anio']}",
                      "Real acumulado + proyección hasta diciembre", COLOR_IE)

    chart_data = CategoryChartData()
    chart_data.categories = datos['grafico_meses']
    chart_data.add_series('Real', datos['grafico_real'])
    chart_data.add_series('Proyección (sin proc. admin.)', datos['grafico_proy_sin'])
    if datos['total_admin'] > 0:
        chart_data.add_series('Proyección (con proc. admin.)', datos['grafico_proy_con'])
    chart_data.add_series('Meta', [datos['meta']] * len(datos['grafico_meses']))

    x, y, cx, cy = Emu(609600), Emu(1341120), Emu(10972800), Emu(4632960)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)
    chart = gf.chart
    _estilizar_chart(chart)
    colores = [COLOR_TEXTO, COLOR_EM, COLOR_CTM, COLOR_MUTED2]
    _colorear_series(chart, colores)

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# Helper — tabla nativa de casos (encabezado + filas alternadas)
# ---------------------------------------------------------
def _tabla_casos(slide, x, y, cx, cy, columnas, filas, color_header_bg):
    """columnas: lista de (titulo, ancho_emu). filas: lista de tuplas de texto, mismo orden que columnas."""
    n_filas = len(filas) + 1
    n_cols = len(columnas)
    gf = slide.shapes.add_table(n_filas, n_cols, x, y, cx, cy)
    table = gf.table
    table.first_row = False
    table.horz_banding = False

    for c, (_, ancho) in enumerate(columnas):
        table.columns[c].width = Emu(ancho)
    alto_fila = Emu(int(cy / n_filas))
    for r in range(n_filas):
        table.rows[r].height = alto_fila

    def _celda(r, c, texto_valor, fondo_hex, color_hex, bold, align):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(fondo_hex)
        cell.margin_left = cell.margin_right = Emu(60960)
        cell.margin_top = cell.margin_bottom = Emu(12192)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(texto_valor)
        run.font.size = Pt(9.5 if r == 0 else 9)
        run.font.bold = bold
        run.font.name = FUENTE
        run.font.color.rgb = _rgb(color_hex)

    for c, (titulo, _) in enumerate(columnas):
        _celda(0, c, titulo, color_header_bg, "FFFFFF", True,
               PP_ALIGN.LEFT if c < 2 else PP_ALIGN.CENTER)

    for r, fila in enumerate(filas, start=1):
        fondo = COLOR_BG if r % 2 == 1 else COLOR_CARD_BG
        for c, valor in enumerate(fila):
            _celda(r, c, valor, fondo, COLOR_TEXTO, False,
                   PP_ALIGN.LEFT if c < 2 else PP_ALIGN.CENTER)
    return gf


# ---------------------------------------------------------
# SLIDES — TOP 10 POR PROBABILIDAD DE CIERRE (ordenados por pérdida bruta)
# ---------------------------------------------------------
_COLUMNAS_TOP10 = [
    ("N° Caso", 950000), ("Nickname", 2350000), ("División", 750000),
    ("Ajustador", 1450000), ("Pérdida Bruta", 1750000), ("Prob.", 979600),
]


def slide_top_probabilidad(prs, datos, clave, titulo, subtitulo, color_accent):
    slide = nueva_slide(prs)
    encabezado_slide(slide, titulo, subtitulo, color_accent)

    casos = datos.get(clave, [])
    if not casos:
        texto(slide, MARGEN_X, Emu(2438400), CONTENIDO_W, Emu(609600),
              "No hay casos que cumplan este criterio.", 16.0, COLOR_MUTED)
    else:
        filas = [
            (
                c['caso'], c['nickname'], c['division'], c['ajustador'],
                f"{c['perdida']:,.0f} {c['divisa']}".replace(",", "."),
                c['probabilidad'],
            )
            for c in casos
        ]
        _tabla_casos(slide, MARGEN_X, Emu(1341120), CONTENIDO_W, Emu(4800000),
                     _COLUMNAS_TOP10, filas, color_accent)

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE — FACTURACIÓN MENSUAL POR DIVISIÓN (barras)
# ---------------------------------------------------------
def slide_facturacion_mensual(prs, datos):
    slide = nueva_slide(prs)
    encabezado_slide(slide, f"Facturación Mensual por División — {datos['nombre_mes_inicio']} a {datos['nombre_mes_corte']} {datos['anio']}",
                      "Real, por mes cerrado", COLOR_IE)

    chart_data = CategoryChartData()
    chart_data.categories = datos['meses_mensual']
    chart_data.add_series('Equipo Móvil', datos['valores_em_mensual'])
    chart_data.add_series('Ingeniería y Energía', datos['valores_ie_mensual'])

    x, y, cx, cy = Emu(609600), Emu(1341120), Emu(10972800), Emu(4632960)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = gf.chart
    _estilizar_chart(chart)
    _colorear_series(chart, [COLOR_EM, COLOR_IE])

    pie_de_pagina(slide, datos['fuente_texto'])
    return slide


# ---------------------------------------------------------
# SLIDE — CIERRE (clonada del formato corporativo JPV)
# ---------------------------------------------------------
def slide_cierre(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    plantilla = _cargar_plantilla_jpv()
    if plantilla is not None and len(plantilla.slides) > 1:
        clonar_shapes_de_plantilla(slide, plantilla.slides[1])
    else:
        # Respaldo si falta el asset de la plantilla (assets/plantilla_jpv.pptx).
        fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        fondo.fill.solid()
        fondo.fill.fore_color.rgb = _rgb(COLOR_BAR)
        _sin_borde(fondo)
        fondo.shadow.inherit = False
        texto(slide, MARGEN_X, Emu(3000000), CONTENIDO_W, Emu(800000),
              "Gracias", 32, "FFFFFF", bold=True, align=PP_ALIGN.CENTER, fuente=FUENTE_TITULO)
    return slide


def generar_pptx_forecast(datos):
    """Genera la presentación completa y devuelve los bytes del archivo .pptx."""
    import io as _io
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_portada(prs, datos)
    slide_agenda(prs, datos)
    slide_titulo_seccion(prs, datos)
    slide_linea_base_real(prs, datos)
    slide_proyeccion_pipeline(prs, datos)
    if datos['total_admin'] > 0:
        slide_casos_admin(prs, datos)
    slide_forecast_total(prs, datos)
    slide_desviacion(prs, datos)
    slide_evolucion(prs, datos)
    slide_facturacion_mensual(prs, datos)
    slide_top_probabilidad(prs, datos, 'top10_100', "Top 10 — Casos con Probabilidad de Cierre 100%",
                            "Ordenados por Pérdida Bruta (mayor a menor)", COLOR_VERDE)
    slide_top_probabilidad(prs, datos, 'top10_75', "Top 10 — Casos con Probabilidad de Cierre 75%",
                            "Ordenados por Pérdida Bruta (mayor a menor)", COLOR_IE)
    slide_top_probabilidad(prs, datos, 'top10_menor50', "Top 10 — Casos con Probabilidad de Cierre < 50%",
                            "Ordenados por Pérdida Bruta (mayor a menor)", COLOR_ROJO)
    slide_cierre(prs, datos)

    buffer = _io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
